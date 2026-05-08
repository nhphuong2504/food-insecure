"""Generate report/slides.pptx -- a story-driven 14-slide deck.

Design philosophy follows the data-storytelling skill:
  Hook -> Context -> Discovery -> Climax -> Recommendation -> Action.

Every slide leads with an insight headline (the "so what"), supports it
with one curated visual, and closes with a take-away strip. Three colors,
two type sizes per slide, no chart legends without a story behind them.
"""

from __future__ import annotations

from pathlib import Path
from PIL import Image as PILImage

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


PROJECT_ROOT = Path(__file__).resolve().parents[1]
FIG = PROJECT_ROOT / "report" / "figures"
OUT = PROJECT_ROOT / "report" / "slides.pptx"


# ----------------------------------------------------------------------
# design system
# ----------------------------------------------------------------------

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY    = RGBColor(0x1B, 0x3A, 0x5B)   # slide title bar, primary accents
ORANGE  = RGBColor(0xE7, 0x6F, 0x51)   # callout numbers, "so what" emphasis
TEAL    = RGBColor(0x2A, 0x9D, 0x8F)   # positive / supporting
CHAR    = RGBColor(0x2B, 0x2B, 0x2B)   # body text
GREY_L  = RGBColor(0xF2, 0xF2, 0xF2)   # bands, light fills
GREY_M  = RGBColor(0xBC, 0xBC, 0xBC)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"


# ----------------------------------------------------------------------
# primitives
# ----------------------------------------------------------------------

def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background() if line is None else None
    if line is None:
        s.line.fill.background()
    return s


def add_text(slide, x, y, w, h, text, *,
             font_size=18, bold=False, color=CHAR, align="left",
             anchor="top", font=FONT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM
    }[anchor]

    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT
    }[align]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, *,
                font_size=16, color=CHAR, accent=NAVY, line_spacing=1.15):
    """Wedge-styled bullets with a colored accent square at the start."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)

        r1 = p.add_run()
        r1.text = "\u25A0  "  # filled square bullet
        r1.font.name = FONT
        r1.font.size = Pt(font_size)
        r1.font.bold = True
        r1.font.color.rgb = accent

        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = color
    return box


def add_image_fit(slide, path: Path, x, y, max_w, max_h):
    """Insert an image scaled to fit a bounding box, centered inside it."""
    if not path.exists():
        return None
    with PILImage.open(path) as im:
        iw, ih = im.size
    ar = ih / iw
    box_w_emu = int(max_w)
    box_h_emu = int(max_h)
    target_h = box_w_emu * ar
    if target_h <= box_h_emu:
        w = box_w_emu
        h = int(target_h)
    else:
        h = box_h_emu
        w = int(box_h_emu / ar)
    cx = int(x) + (box_w_emu - w) // 2
    cy = int(y) + (box_h_emu - h) // 2
    return slide.shapes.add_picture(str(path), cx, cy, width=w, height=h)


# ----------------------------------------------------------------------
# slide chrome
# ----------------------------------------------------------------------

def add_header(slide, slide_title: str, kicker: str = ""):
    """Slim navy bar with white title; small kicker above (eyebrow)."""
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.85), NAVY)
    if kicker:
        add_text(slide, Inches(0.5), Inches(0.10), Inches(8), Inches(0.25),
                 kicker.upper(), font_size=10, bold=True, color=GREY_L,
                 anchor="middle")
    add_text(slide, Inches(0.5), Inches(0.30), Inches(12.3), Inches(0.5),
             slide_title, font_size=22, bold=True, color=WHITE,
             anchor="middle")


def add_footer(slide, slide_no: int, total: int, label: str = "EE 5290 - Final Project"):
    add_rect(slide, Emu(0), Inches(7.20), SLIDE_W, Inches(0.30), GREY_L)
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.26),
             label, font_size=9, color=CHAR, anchor="middle")
    add_text(slide, Inches(11.5), Inches(7.22), Inches(1.3), Inches(0.26),
             f"{slide_no} / {total}", font_size=9, color=CHAR,
             align="right", anchor="middle")


def add_takeaway(slide, text: str, color=ORANGE, y=Inches(6.45)):
    """Bottom take-away strip -- the 'so what' of every slide."""
    add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.55), color)
    add_text(slide, Inches(0.7), y, Inches(11.9), Inches(0.55),
             text, font_size=15, bold=True, color=WHITE, anchor="middle")


# ----------------------------------------------------------------------
# slides
# ----------------------------------------------------------------------

def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# Slide 1 -- Title -----------------------------------------------------
def slide_01_title(prs):
    slide = slide_blank(prs)
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, Emu(0), Inches(2.2), SLIDE_W, Inches(0.05), ORANGE)

    add_text(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             "EE 5290 - DATA ANALYTICS - SPRING 2026",
             font_size=12, bold=True, color=ORANGE)
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12.2), Inches(1.5),
             "Targeting Food Insecurity\nin the Midwest",
             font_size=48, bold=True, color=WHITE)
    add_text(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.7),
             "From 4,826 households to two programs that matter",
             font_size=22, color=GREY_L, italic=True)
    add_text(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             "Final Project  -  Iowa State University  -  May 2026",
             font_size=12, color=GREY_L)


# Slide 2 -- Hook (the headline insight) -------------------------------
def slide_02_hook(prs):
    slide = slide_blank(prs)
    add_text(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.4),
             "THE HEADLINE",
             font_size=12, bold=True, color=ORANGE)
    add_text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(2.0),
             '"1 in 4 Midwest households is food insecure.\n'
             'A simple model finds 2 of every 3 it flags --\n'
             'and points to the two programs that solve 88% of the need."',
             font_size=30, bold=True, color=NAVY, italic=True)

    stat_y = Inches(4.3)
    stat_h = Inches(2.0)
    col_w  = Inches(4.0)
    gap    = Inches(0.15)
    x0     = Inches(0.5)

    def stat_card(x, big, small, color):
        add_rect(slide, x, stat_y, col_w, stat_h, GREY_L)
        add_rect(slide, x, stat_y, col_w, Inches(0.10), color)
        add_text(slide, x, stat_y + Inches(0.30), col_w, Inches(1.10),
                 big, font_size=58, bold=True, color=color, align="center")
        add_text(slide, x, stat_y + Inches(1.45), col_w, Inches(0.50),
                 small, font_size=13, color=CHAR, align="center", anchor="middle")

    stat_card(x0,                                   "27.85%",
              "of households flagged food-insecure", NAVY)
    stat_card(x0 + col_w + gap,                     "67%",
              "precision among the top-decile risk score", ORANGE)
    stat_card(x0 + 2 * (col_w + gap),               "88%",
              "of those flags belong to just 2 personas", TEAL)

    add_takeaway(slide,
                 "These three numbers carry the whole talk. Everything else explains them.")


# Slide 3 -- The Problem -----------------------------------------------
def slide_03_problem(prs):
    slide = slide_blank(prs)
    add_header(slide, "Demand is rising. Capacity isn't.", kicker="The Problem")

    add_text(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "Midwest food banks must decide  *which*  households to prioritize and  "
             "*what kind*  of program to fund -- with no slack on either side.",
             font_size=18, color=CHAR, italic=True)

    # Two question panels
    panel_y = Inches(2.5)
    panel_h = Inches(3.5)
    panel_w = Inches(5.9)
    gap     = Inches(0.5)
    x_left  = Inches(0.5)
    x_right = x_left + panel_w + gap

    add_rect(slide, x_left, panel_y, panel_w, panel_h, GREY_L)
    add_rect(slide, x_left, panel_y, Inches(0.18), panel_h, NAVY)
    add_text(slide, x_left + Inches(0.4), panel_y + Inches(0.25),
             panel_w - Inches(0.5), Inches(0.5),
             "Question 1  -  Prediction",
             font_size=18, bold=True, color=NAVY)
    add_text(slide, x_left + Inches(0.4), panel_y + Inches(0.95),
             panel_w - Inches(0.5), Inches(2.2),
             "Can a model identify food-insecure households "
             "accurately enough to drive triage?\n\n"
             "If yes -> we can rank households and target outreach.",
             font_size=15, color=CHAR)

    add_rect(slide, x_right, panel_y, panel_w, panel_h, GREY_L)
    add_rect(slide, x_right, panel_y, Inches(0.18), panel_h, ORANGE)
    add_text(slide, x_right + Inches(0.4), panel_y + Inches(0.25),
             panel_w - Inches(0.5), Inches(0.5),
             "Question 2  -  Segmentation",
             font_size=18, bold=True, color=ORANGE)
    add_text(slide, x_right + Inches(0.4), panel_y + Inches(0.95),
             panel_w - Inches(0.5), Inches(2.2),
             "Among food-insecure households, what  *kinds*  exist?\n\n"
             "If we know -> we can design programs that fit, "
             "not one-size-fits-all distribution.",
             font_size=15, color=CHAR)

    add_takeaway(slide,
                 "Two questions, two tracks, one synthesis. That's the structure of this talk.")


# Slide 4 -- The Data --------------------------------------------------
def slide_04_data(prs):
    slide = slide_blank(prs)
    add_header(slide, "4,826 households, 38 features, one stubborn imbalance",
               kicker="The Data")

    # Left: facts
    add_text(slide, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.45),
             "What we have", font_size=16, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.5), Inches(1.65), Inches(6.0), Inches(4.5), [
        "FoodAPS-derived household extract; one row per household",
        "Outcome:  food_insecure_flag_adult  (binary)",
        "27.85% positive -> moderate but real class imbalance",
        "Sentinel codes (-996/-997/-998) -> NaN; caraccess dropped (92% missing)",
        "Distance variables logged; one-hot encoding for nominals",
        "Stratified 80/20 split; 5-fold CV on train only -- test set untouched until the end",
    ], font_size=14)

    # Right: the class-balance figure
    add_image_fit(slide, FIG / "01_class_balance.png",
                  Inches(7.0), Inches(1.10), Inches(5.9), Inches(5.0))

    add_takeaway(slide,
                 "Accuracy will mislead. We will judge models by PR-AUC throughout.")


# Slide 5 -- EDA insight -----------------------------------------------
def slide_05_eda(prs):
    slide = slide_blank(prs)
    add_header(slide, "Poverty nearly doubles the risk. Pantry use triples it.",
               kicker="What the data says before any model")

    add_image_fit(slide, FIG / "02_marginal_rates.png",
                  Inches(0.5), Inches(1.10), Inches(12.3), Inches(5.1))

    add_takeaway(slide,
                 "Bivariate associations point us at the right features -- but they don't control for "
                 "confounding. That's why we go to a model next.")


# Slide 6 -- Approach (methods on a single slide) ----------------------
def slide_06_approach(prs):
    slide = slide_blank(prs)
    add_header(slide, "Two tracks, both honest about uncertainty",
               kicker="The Approach")

    # Track 1
    panel_y = Inches(1.2)
    panel_h = Inches(5.0)
    panel_w = Inches(6.0)
    x_left  = Inches(0.5)
    x_right = Inches(6.85)

    add_rect(slide, x_left, panel_y, panel_w, panel_h, GREY_L)
    add_rect(slide, x_left, panel_y, panel_w, Inches(0.6), NAVY)
    add_text(slide, x_left, panel_y, panel_w, Inches(0.6),
             "TRACK 1  -  Predictive Modelling",
             font_size=15, bold=True, color=WHITE, align="center", anchor="middle")
    add_bullets(slide, x_left + Inches(0.3), panel_y + Inches(0.85),
                panel_w - Inches(0.6), Inches(4.0), [
        "Three classifiers:  Logistic Regression,  SVM-RBF,  Random Forest",
        "Three imbalance strategies:  none,  balanced weights,  F1-tuned threshold",
        "9-cell results matrix on a single 80/20 hold-out",
        "Tuning metric:  average precision (PR-AUC)",
        "Threshold tuned on TRAIN out-of-fold predictions only",
    ], font_size=13, accent=NAVY)

    add_rect(slide, x_right, panel_y, panel_w, panel_h, GREY_L)
    add_rect(slide, x_right, panel_y, panel_w, Inches(0.6), ORANGE)
    add_text(slide, x_right, panel_y, panel_w, Inches(0.6),
             "TRACK 2  -  Segmentation",
             font_size=15, bold=True, color=WHITE, align="center", anchor="middle")
    add_bullets(slide, x_right + Inches(0.3), panel_y + Inches(0.85),
                panel_w - Inches(0.6), Inches(4.0), [
        "Filter to the 1,344 food-insecure households",
        "Standardize -> PCA (24 PCs explain 80% of variance)",
        "K-Means sweep  k=2..8;  choose by silhouette + elbow",
        "Ran TWICE:  with geography (k=2 rural/urban),  without (k=5 personas)",
        "Stability check:  8 random inits,  mean ARI > 0.9",
    ], font_size=13, accent=ORANGE)

    add_takeaway(slide,
                 "Track 1 ranks households. Track 2 names them. Synthesis joins them at the end.")


# Slide 7 -- Track 1 result --------------------------------------------
def slide_07_track1_result(prs):
    slide = slide_blank(prs)
    add_header(slide, "Logistic regression wins.  PR-AUC 0.568,  ROC-AUC 0.78.",
               kicker="Track 1 - Result")

    # left: the 3-row summary table (best per family)
    add_text(slide, Inches(0.5), Inches(1.1), Inches(7.0), Inches(0.35),
             "Best operating point per model family",
             font_size=13, bold=True, color=NAVY)

    headers = ["Model (best variant)", "PR-AUC", "ROC-AUC", "Prec.", "Recall", "F1"]
    rows = [
        ["Logistic Regression  (balanced)",   "0.568", "0.780", "0.463", "0.762", "0.576"],
        ["SVM-RBF  (F1-tuned threshold)",     "0.529", "0.770", "0.425", "0.747", "0.542"],
        ["Random Forest  (F1-tuned threshold)","0.521", "0.760", "0.453", "0.792", "0.577"],
    ]

    table_left = Inches(0.5)
    table_top  = Inches(1.55)
    col_widths = [Inches(2.6), Inches(0.8), Inches(0.85), Inches(0.85), Inches(0.85), Inches(0.85)]
    table_w_total = sum([cw for cw in col_widths], Inches(0))
    row_h = Inches(0.45)

    # header row
    cx = table_left
    for i, h in enumerate(headers):
        add_rect(slide, cx, table_top, col_widths[i], row_h, NAVY)
        add_text(slide, cx, table_top, col_widths[i], row_h, h,
                 font_size=12, bold=True, color=WHITE, align="center", anchor="middle")
        cx = cx + col_widths[i]

    # data rows -- highlight LR row (winner)
    for ri, row in enumerate(rows):
        cy = table_top + row_h * (ri + 1)
        bg = (ORANGE if ri == 0 else GREY_L)
        cx = table_left
        for ci, val in enumerate(row):
            add_rect(slide, cx, cy, col_widths[ci], row_h,
                     bg if ri == 0 else (WHITE if ri % 2 else GREY_L))
            txt_color = WHITE if ri == 0 else CHAR
            bold = (ri == 0)
            align = "left" if ci == 0 else "center"
            add_text(slide, cx + Inches(0.08), cy, col_widths[ci] - Inches(0.16), row_h,
                     val, font_size=12, bold=bold, color=txt_color,
                     align=align, anchor="middle")
            cx = cx + col_widths[ci]

    # right: three story bullets
    add_text(slide, Inches(8.4), Inches(1.1), Inches(4.5), Inches(0.35),
             "Three takeaways", font_size=13, bold=True, color=NAVY)
    add_bullets(slide, Inches(8.4), Inches(1.55), Inches(4.5), Inches(4.4), [
        "LR beats RF and SVM by a clear ~0.04 PR-AUC -- the signal is mostly linear after feature engineering.",
        "Threshold tuning roughly equals balanced weights -- the *threshold* matters more than the *weighting*.",
        "If we demand 70% precision, we recover only ~25% of the food-insecure -- a real ceiling.",
    ], font_size=13, accent=ORANGE)

    add_takeaway(slide,
                 "We chose interpretability without giving anything up. That matters for what comes next.")


# Slide 8 -- The Tradeoff ----------------------------------------------
def slide_08_tradeoff(prs):
    slide = slide_blank(prs)
    add_header(slide,
               "Targeting the top 10% catches 1 in 4 food-insecure households at 67% precision.",
               kicker="Track 1 - The Operating Point")

    add_image_fit(slide, FIG / "05_pr_curves.png",
                  Inches(0.5), Inches(1.1), Inches(6.2), Inches(5.1))
    add_image_fit(slide, FIG / "06_calibration.png",
                  Inches(7.0), Inches(1.1), Inches(5.9), Inches(5.1))

    add_takeaway(slide,
                 "LR dominates in the high-recall regime that matters for triage -- and stays well calibrated.")


# Slide 9 -- What drives it --------------------------------------------
def slide_09_drivers(prs):
    slide = slide_blank(prs)
    add_header(slide,
               "Poverty band, pantry use, education -- the same story two ways.",
               kicker="Track 1 - What the Models Look At")

    add_image_fit(slide, FIG / "08_feature_importance.png",
                  Inches(0.5), Inches(1.05), Inches(12.3), Inches(5.0))

    add_takeaway(slide,
                 "Linear coefficients (left) and tree importances (right) agree on the top features. "
                 "That's the strongest evidence the signal is real.")


# Slide 10 -- Track 2 setup --------------------------------------------
def slide_10_track2_setup(prs):
    slide = slide_blank(prs)
    add_header(slide, "Five personas hide inside the food-insecure population.",
               kicker="Track 2 - The Discovery")

    # left: explanation
    add_bullets(slide, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.1), [
        "1,344 food-insecure households -- the segmentation universe.",
        "Standardize, PCA, sweep K-Means k = 2..8.",
        "With geography:  silhouette winner is k=2.  Just rural vs urban.",
        "Drop distance/region features  ->  silhouette + elbow agree on  k = 5.",
        "Eight random initializations agree (mean ARI > 0.9):  the personas are stable.",
    ], font_size=14, accent=ORANGE)

    # right: the silhouette/elbow chart
    add_image_fit(slide, FIG / "10_kmeans_sweep.png",
                  Inches(6.85), Inches(1.05), Inches(6.0), Inches(5.0))

    add_takeaway(slide,
                 "Geography is the dominant axis -- but if we set it aside, "
                 "five operationally meaningful personas appear.")


# Slide 11 -- Meet the personas ----------------------------------------
def slide_11_personas_table(prs):
    slide = slide_blank(prs)
    add_header(slide, "Meet the five personas", kicker="Track 2 - The Personas")

    # left: heatmap (small)
    add_image_fit(slide, FIG / "14_persona_heatmap.png",
                  Inches(0.4), Inches(1.05), Inches(6.4), Inches(5.1))

    # right: persona table
    table_left = Inches(7.0)
    table_top  = Inches(1.05)
    col_w = [Inches(0.55), Inches(2.7), Inches(0.55), Inches(2.4)]
    row_h = Inches(0.92)
    rows = [
        ["P0", "Working families just above poverty", "20%",
         "100% employed; pov. ratio 1.43"],
        ["P1", "Non-working low-income singles", "29%",
         "Pantry use 28% (3x avg)"],
        ["P2", "Large low-income families with children", "26%",
         "Mean HH 5.0; 2.6 kids; pov. 0.85"],
        ["P3", "Working educated, near boundary", "16%",
         "Pov. ratio 3.05 (above poverty)"],
        ["P4", "Elderly low-income, high SNAP", "9%",
         "Mean head age 70; 42% on SNAP"],
    ]
    persona_colors = [TEAL, ORANGE, ORANGE, NAVY, NAVY]
    # header
    headers = ["#", "Persona name", "Share", "Headline feature"]
    cx = table_left
    for i, h in enumerate(headers):
        add_rect(slide, cx, table_top, col_w[i], Inches(0.4), NAVY)
        add_text(slide, cx + Inches(0.05), table_top, col_w[i] - Inches(0.1),
                 Inches(0.4), h, font_size=11, bold=True, color=WHITE,
                 align="center" if i != 1 else "left", anchor="middle")
        cx = cx + col_w[i]

    for ri, row in enumerate(rows):
        cy = table_top + Inches(0.4) + row_h * ri
        bg = WHITE if ri % 2 == 0 else GREY_L
        cx = table_left
        # P# colored cell
        add_rect(slide, cx, cy, col_w[0], row_h, persona_colors[ri])
        add_text(slide, cx, cy, col_w[0], row_h, row[0],
                 font_size=14, bold=True, color=WHITE, align="center", anchor="middle")
        cx = cx + col_w[0]
        for ci in range(1, 4):
            add_rect(slide, cx, cy, col_w[ci], row_h, bg)
            align = "left" if ci == 1 else ("center" if ci == 2 else "left")
            font_sz = 11
            add_text(slide, cx + Inches(0.08), cy, col_w[ci] - Inches(0.16), row_h,
                     row[ci], font_size=font_sz, color=CHAR,
                     align=align, anchor="middle")
            cx = cx + col_w[ci]

    add_takeaway(slide,
                 "P1 and P2 together are 55% of the food-insecure -- and they look completely different.")


# Slide 12 -- The CLIMAX (synthesis) -----------------------------------
def slide_12_climax(prs):
    slide = slide_blank(prs)
    add_header(slide,
               "88% of high-risk flags belong to just two personas.",
               kicker="The Synthesis - The Climax")

    # Left big number
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(4.0), Inches(4.9), GREY_L)
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(4.0), Inches(0.10), ORANGE)
    add_text(slide, Inches(0.5), Inches(1.5), Inches(4.0), Inches(1.0),
             "TOP-DECILE FLAGS",
             font_size=13, bold=True, color=NAVY, align="center", anchor="middle")
    add_text(slide, Inches(0.5), Inches(2.4), Inches(4.0), Inches(1.5),
             "88%",
             font_size=88, bold=True, color=ORANGE, align="center", anchor="middle")
    add_text(slide, Inches(0.5), Inches(4.0), Inches(4.0), Inches(0.7),
             "are P1 (45%) + P2 (43%)",
             font_size=18, bold=True, color=CHAR, align="center", anchor="middle")
    add_text(slide, Inches(0.6), Inches(4.8), Inches(3.8), Inches(1.1),
             "Non-working singles  +  large families with children. "
             "Two personas. Two distinct programs.",
             font_size=12, color=CHAR, align="center", italic=True)

    # Right chart -- persona mix
    add_image_fit(slide, FIG / "15_persona_mix_at_cutoffs.png",
                  Inches(4.8), Inches(1.15), Inches(8.2), Inches(4.95))

    add_takeaway(slide,
                 "You don't need five programs. You need two -- and the model tells you exactly who they are.")


# Slide 13 -- Recommendation -------------------------------------------
def slide_13_recommendation(prs):
    slide = slide_blank(prs)
    add_header(slide, "Two programs, two precise targeting strategies",
               kicker="The Recommendation")

    # Two program cards on top
    card_y = Inches(1.10)
    card_h = Inches(2.7)
    card_w = Inches(6.0)
    x_left  = Inches(0.5)
    x_right = Inches(6.85)

    def program_card(x, color, title, share, share_lbl, bullets):
        add_rect(slide, x, card_y, card_w, card_h, GREY_L)
        add_rect(slide, x, card_y, Inches(0.18), card_h, color)
        add_text(slide, x + Inches(0.4), card_y + Inches(0.15),
                 card_w - Inches(2.5), Inches(0.45),
                 title, font_size=15, bold=True, color=color)
        add_text(slide, x + card_w - Inches(2.0), card_y + Inches(0.10),
                 Inches(1.8), Inches(0.6),
                 share, font_size=32, bold=True, color=color,
                 align="right", anchor="middle")
        add_text(slide, x + card_w - Inches(2.0), card_y + Inches(0.65),
                 Inches(1.8), Inches(0.3),
                 share_lbl, font_size=10, color=CHAR,
                 align="right", anchor="middle")
        add_bullets(slide, x + Inches(0.4), card_y + Inches(1.10),
                    card_w - Inches(0.8), card_h - Inches(1.2),
                    bullets, font_size=12, accent=color)

    program_card(x_left, ORANGE,
        "PROGRAM A  -  Pantry partnerships for low-income singles",
        "45%", "of top-decile flags",
        ["Targets P1:  non-working low-income singles, pantry use 28% (3x average)",
         "These households already self-select into pantries -- meet them where they are",
         "Operational lift:  partner with existing pantries; supplemental distribution"])

    program_card(x_right, NAVY,
        "PROGRAM B  -  Family bulk packs and school-meal coordination",
        "43%", "of top-decile flags",
        ["Targets P2:  mean household size 5.0, 2.6 kids, poverty ratio 0.85",
         "School-meal eligibility is highest here -- coordinate with school districts",
         "Operational lift:  bulk-pack distribution sized for families"])

    # Equity caveat panel
    eq_y = Inches(4.0)
    add_rect(slide, Inches(0.5), eq_y, Inches(12.3), Inches(2.3), GREY_L)
    add_rect(slide, Inches(0.5), eq_y, Inches(12.3), Inches(0.10), TEAL)
    add_text(slide, Inches(0.7), eq_y + Inches(0.15), Inches(12.0), Inches(0.45),
             "Equity caveat:  who the model misses",
             font_size=15, bold=True, color=TEAL)
    add_bullets(slide, Inches(0.7), eq_y + Inches(0.65),
                Inches(12.0), Inches(1.6), [
        "P3  (working educated, just above poverty):  37% of population, <0.5% of top-decile flags. "
        "Use qualitative outreach -- they're invisible to income-eligibility programs.",
        "P4  (elderly on SNAP):  9% of food insecure, under-flagged. "
        "Monitor via existing partnerships; do not disengage on a low score.",
        "All findings are associations, not causal claims.  Survey weights and county-level data are next steps.",
    ], font_size=12, accent=TEAL)

    add_takeaway(slide,
                 "The model is a triage tool, not a verdict. Pair it with the personas to design programs.",
                 color=NAVY)


# Slide 14 -- The Ask --------------------------------------------------
def slide_14_ask(prs):
    slide = slide_blank(prs)
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, Emu(0), Inches(2.2), SLIDE_W, Inches(0.05), ORANGE)

    add_text(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             "REPRODUCIBLE  -  AUDITABLE  -  READY",
             font_size=12, bold=True, color=ORANGE)
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12.2), Inches(1.2),
             "Thank you.  Questions?",
             font_size=44, bold=True, color=WHITE)
    add_text(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
             "All 16 figures, the 9-cell results table, and the persona profiles "
             "are reproducible with one command:",
             font_size=15, color=GREY_L, italic=True)
    add_rect(slide, Inches(0.6), Inches(4.95), Inches(12.0), Inches(0.65), CHAR)
    add_text(slide, Inches(0.8), Inches(4.95), Inches(11.6), Inches(0.65),
             "$  pip install -r requirements.txt  &&  jupyter nbconvert --execute notebooks/*.ipynb",
             font_size=14, color=WHITE, font="Consolas", anchor="middle")
    add_text(slide, Inches(0.6), Inches(6.0), Inches(12.0), Inches(0.4),
             "Repository:  notebooks/  +  src/  +  report/   -   "
             "random_state = 42  throughout.",
             font_size=12, color=GREY_L)


# ----------------------------------------------------------------------
# driver
# ----------------------------------------------------------------------

SLIDES = [
    slide_01_title,
    slide_02_hook,
    slide_03_problem,
    slide_04_data,
    slide_05_eda,
    slide_06_approach,
    slide_07_track1_result,
    slide_08_tradeoff,
    slide_09_drivers,
    slide_10_track2_setup,
    slide_11_personas_table,
    slide_12_climax,
    slide_13_recommendation,
    slide_14_ask,
]


def render() -> None:
    prs = make_presentation()
    total = len(SLIDES)
    for idx, slide_fn in enumerate(SLIDES, start=1):
        slide_fn(prs)
        # Footer on all non-cover slides
        if idx not in (1, total):
            add_footer(prs.slides[idx - 1], idx, total)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}  ({total} slides)")


if __name__ == "__main__":
    render()
